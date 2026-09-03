"""Per-format document extractors for the parsing sidecar (ADR-0006, SPEC-05 §2).

Each ``parse_*`` turns the raw bytes of one heavy format (PDF, DOCX, PPTX, XLSX)
into the same normalised shape the Go parsers emit (``internal/ingest/parse``):

    {"title": str, "blocks": [ {type, level?, text?, rows?}, ... ]}

so a Go parser and this sidecar are interchangeable to the downstream chunker
(SPEC-05 §3). Block types are the closed set ``heading | paragraph | table |
list | code``. Tables carry both structured ``rows`` (header row first) and a GFM
markdown ``text`` — identical rendering to the Go ``markdownTable`` so hashing is
stable regardless of which producer parsed the document.

Pure functions: bytes in, dict out. No network, no filesystem, no globals.
"""

from __future__ import annotations

import io
from collections import Counter

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from docx.document import Document as DocxDocumentType
from docx.oxml.ns import qn
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph
from openpyxl import load_workbook
from pptx import Presentation

# A worksheet or a giant PDF could carry unbounded rows; cap so one pathological
# document cannot blow up memory or the response. ponytail: a fixed cap, not a
# streaming parse — upgrade path is pagination if huge sheets ever matter.
MAX_TABLE_ROWS = 2000


def block(type_: str, *, level: int = 0, text: str = "", rows=None) -> dict:
    """Build a block dict, omitting zero/empty fields to match the Go JSON
    (``omitempty`` on level/text/rows)."""
    b: dict = {"type": type_}
    if level:
        b["level"] = level
    if text:
        b["text"] = text
    if rows:
        b["rows"] = rows
    return b


def markdown_table(rows) -> str:
    """Render rows (header first) as a GFM markdown table, escaping pipes. Mirrors
    the Go ``markdownTable`` byte-for-byte."""
    rows = [r for r in rows]
    if not rows:
        return ""
    cols = max((len(r) for r in rows), default=0)
    if cols == 0:
        return ""

    def render(cells) -> str:
        out = ["|"]
        for c in range(cols):
            cell = str(cells[c]).replace("|", "\\|") if c < len(cells) else ""
            out.append(" " + cell + " |")
        return "".join(out)

    lines = [render(rows[0]), "|" + " --- |" * cols]
    lines.extend(render(r) for r in rows[1:])
    return "\n".join(lines)


def _table_block(rows) -> dict | None:
    """A table block from rows, or None when there is nothing to render. Rows are
    stringified and truncated to MAX_TABLE_ROWS."""
    clean = [[("" if v is None else str(v)) for v in r] for r in rows[:MAX_TABLE_ROWS]]
    clean = [r for r in clean if any(cell.strip() for cell in r)]
    if not clean:
        return None
    return block("table", rows=clean, text=markdown_table(clean))


# --- DOCX ---------------------------------------------------------------------


def parse_docx(data: bytes) -> dict:
    doc: DocxDocumentType = DocxDocument(io.BytesIO(data))
    title = (doc.core_properties.title or "").strip()
    blocks: list[dict] = []

    # Walk the body in document order so paragraphs and tables interleave the way
    # the author wrote them (python-docx exposes them as separate collections).
    body = doc.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            para = DocxParagraph(child, doc)
            text = para.text.strip()
            if not text:
                continue
            level = _docx_heading_level(para.style.name if para.style else "")
            if level:
                blocks.append(block("heading", level=level, text=text))
                if not title and level == 1:
                    title = text
            else:
                blocks.append(block("paragraph", text=text))
        elif child.tag == qn("w:tbl"):
            tbl = DocxTable(child, doc)
            rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
            tb = _table_block(rows)
            if tb:
                blocks.append(tb)

    if not title:
        title = _first_heading(blocks)
    return {"title": title, "blocks": blocks}


def _docx_heading_level(style_name: str) -> int:
    """Map a paragraph style ("Heading 2", "Title") to a 1-6 heading level, else 0."""
    name = (style_name or "").strip().lower()
    if name == "title":
        return 1
    if name.startswith("heading "):
        try:
            return max(1, min(6, int(name[len("heading ") :])))
        except ValueError:
            return 0
    return 0


# --- PPTX ---------------------------------------------------------------------


def parse_pptx(data: bytes) -> dict:
    prs = Presentation(io.BytesIO(data))
    blocks: list[dict] = []
    title = ""

    for slide in prs.slides:
        slide_title = ""
        if slide.shapes.title is not None:
            slide_title = (slide.shapes.title.text or "").strip()
        if slide_title:
            blocks.append(block("heading", level=1, text=slide_title))
            if not title:
                title = slide_title

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_table:
                rows = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ]
                tb = _table_block(rows)
                if tb:
                    blocks.append(tb)
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        blocks.append(block("paragraph", text=text))

    return {"title": title, "blocks": blocks}


# --- XLSX ---------------------------------------------------------------------


def parse_xlsx(data: bytes) -> dict:
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    blocks: list[dict] = []
    for ws in wb.worksheets:
        rows = [list(r) for r in ws.iter_rows(values_only=True)]
        tb = _table_block(rows)
        if tb is None:
            continue
        blocks.append(block("heading", level=2, text=ws.title))
        blocks.append(tb)
    wb.close()
    # A workbook carries no document title; it comes from source metadata upstream.
    return {"title": "", "blocks": blocks}


# --- PDF ----------------------------------------------------------------------


def parse_pdf(data: bytes) -> dict:
    doc = fitz.open(stream=data, filetype="pdf")
    title = (doc.metadata.get("title") or "").strip() if doc.metadata else ""
    blocks: list[dict] = []

    body_size = _pdf_body_size(doc)
    for page in doc:
        _pdf_page_blocks(page, body_size, blocks)

    if not title:
        title = _first_heading(blocks)
    doc.close()
    return {"title": title, "blocks": blocks}


def _pdf_body_size(doc) -> float:
    """The most common span font size across the document — the body text size,
    the baseline heading detection compares against."""
    sizes: Counter = Counter()
    for page in doc:
        for b in page.get_text("dict")["blocks"]:
            for line in b.get("lines", []):
                for span in line.get("spans", []):
                    if span["text"].strip():
                        sizes[round(span["size"], 1)] += len(span["text"])
    return sizes.most_common(1)[0][0] if sizes else 0.0


def _pdf_page_blocks(page, body_size: float, blocks: list[dict]) -> None:
    """Extract a page's text blocks in reading order, classifying visibly-larger
    lines as headings.

    ponytail: heading detection is font-size ranking, not a trained layout model,
    and tables inside a PDF are appended after the page text rather than placed
    inline (find_tables gives bboxes but interleaving them by position is fiddly).
    Ceiling: a heading may be mis-levelled and a table may sit slightly out of
    order — never lost. Upgrade path: Docling/Unstructured for true layout.
    """
    table_bboxes = []
    try:
        found = page.find_tables()
        for t in found.tables:
            rows = t.extract()
            tb = _table_block(rows)
            if tb:
                table_bboxes.append(fitz.Rect(t.bbox))
    except Exception:  # find_tables is best-effort; never fail the page for it.
        found = None

    for b in sorted(page.get_text("dict")["blocks"], key=lambda b: (b["bbox"][1], b["bbox"][0])):
        if "lines" not in b:
            continue
        brect = fitz.Rect(b["bbox"])
        if any(brect.intersects(tb) for tb in table_bboxes):
            continue  # text captured by a table; emit the table instead (below).
        for line in b["lines"]:
            text = "".join(span["text"] for span in line["spans"]).strip()
            if not text:
                continue
            max_size = max((span["size"] for span in line["spans"]), default=body_size)
            level = _pdf_heading_level(max_size, body_size)
            if level:
                blocks.append(block("heading", level=level, text=text))
            else:
                blocks.append(block("paragraph", text=text))

    # Emit the page's tables after its text (documented ceiling above).
    if found is not None:
        for t in found.tables:
            tb = _table_block(t.extract())
            if tb:
                blocks.append(tb)


def _pdf_heading_level(size: float, body_size: float) -> int:
    """A line's heading level from how much larger it is than body text, else 0."""
    if body_size <= 0 or size < body_size * 1.15:
        return 0
    ratio = size / body_size
    if ratio >= 1.7:
        return 1
    if ratio >= 1.35:
        return 2
    return 3


# --- shared -------------------------------------------------------------------


def _first_heading(blocks) -> str:
    for b in blocks:
        if b["type"] == "heading":
            return b.get("text", "")
    return ""


# MIME → extractor. The keys are the canonical upload MIME types (SPEC-05 §2); the
# app dispatches on the ``mime`` form field, never on sniffing.
DISPATCH = {
    "application/pdf": parse_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": parse_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": parse_pptx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": parse_xlsx,
}
