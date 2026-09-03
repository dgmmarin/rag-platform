"""Generate the six binary parser fixtures under ``testdata/`` (STORY-05.3).

Run with ``python gen_fixtures.py`` (needs the runtime deps in requirements.txt).
The fixtures are committed so the test suite stays hermetic; this script keeps
them regenerable rather than opaque binaries. It uses only the runtime parsing
libraries — no extra dependency for fixture authoring.
"""

from __future__ import annotations

import os

import fitz  # PyMuPDF
from docx import Document as Docx
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

HERE = os.path.join(os.path.dirname(__file__), "testdata")


def _pdf(path: str, title: str, lines: list[tuple[str, int]], pages: int = 1) -> None:
    doc = fitz.open()
    doc.set_metadata({"title": title})
    per = max(1, len(lines) // pages)
    for p in range(pages):
        page = doc.new_page()
        y = 72
        for text, size in lines[p * per : (p + 1) * per] if p < pages - 1 else lines[p * per :]:
            page.insert_text((72, y), text, fontsize=size)
            y += size * 1.7
    doc.save(path)
    doc.close()


def gen_report_pdf() -> None:
    # Headings (large) + body (small): exercises font-size heading detection.
    _pdf(
        os.path.join(HERE, "report.pdf"),
        "Quarterly Report",
        [
            ("Quarterly Report", 24),
            ("Executive Summary", 18),
            ("Revenue grew across every region this quarter.", 11),
            ("Costs were held flat against the prior period.", 11),
            ("Outlook", 18),
            ("We expect continued growth into the next quarter.", 11),
        ],
    )


def gen_article_pdf() -> None:
    _pdf(
        os.path.join(HERE, "article.pdf"),
        "On Retrieval",
        [
            ("On Retrieval", 24),
            ("Why grounding matters", 16),
            ("Grounded answers cite their sources and refuse when unsure.", 11),
            ("Chunking", 16),
            ("Structure-aware chunks keep tables and headings intact.", 11),
        ],
        pages=2,
    )


def gen_memo_docx() -> None:
    d = Docx()
    d.core_properties.title = "Team Memo"
    d.add_heading("Team Memo", level=1)
    d.add_paragraph("Please review the rollout plan below before Friday.")
    d.add_heading("Milestones", level=2)
    t = d.add_table(rows=3, cols=2)
    data = [["Milestone", "Date"], ["Beta", "2026-10-01"], ["GA", "2026-11-15"]]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            t.cell(r, c).text = val
    d.save(os.path.join(HERE, "memo.docx"))


def gen_letter_docx() -> None:
    d = Docx()
    d.core_properties.title = "Offer Letter"
    d.add_heading("Offer Letter", level=1)
    d.add_paragraph("We are pleased to extend the following offer.")
    d.add_heading("Terms", level=2)
    t = d.add_table(rows=3, cols=2)
    data = [["Field", "Value"], ["Role", "Engineer"], ["Start", "2026-10-01"]]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            t.cell(r, c).text = val
    d.save(os.path.join(HERE, "letter.docx"))


def gen_deck_pptx() -> None:
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "Product Overview"
    body = s1.placeholders[1].text_frame
    body.text = "Multi-tenant by design"
    body.add_paragraph().text = "Best-in-class retrieval"

    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Key Metrics"
    tbl = s2.shapes.add_table(3, 2, Inches(1), Inches(1.5), Inches(6), Inches(2)).table
    data = [["Metric", "Value"], ["Tenants", "42"], ["Uptime", "99.9%"]]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val
    prs.save(os.path.join(HERE, "deck.pptx"))


def gen_budget_xlsx() -> None:
    wb = Workbook()
    q1 = wb.active
    q1.title = "Q1"
    for row in [["Item", "Amount"], ["Compute", 1200], ["Storage", 300]]:
        q1.append(row)
    q2 = wb.create_sheet("Q2")
    for row in [["Item", "Amount"], ["Compute", 1500], ["Storage", 350]]:
        q2.append(row)
    wb.save(os.path.join(HERE, "budget.xlsx"))


def main() -> None:
    os.makedirs(HERE, exist_ok=True)
    gen_report_pdf()
    gen_article_pdf()
    gen_memo_docx()
    gen_letter_docx()
    gen_deck_pptx()
    gen_budget_xlsx()
    print(f"wrote fixtures to {HERE}")


if __name__ == "__main__":
    main()
